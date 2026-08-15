"""Build GuideLearn JUMP THAILAND 2026 pitch deck (editable PPTX)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

# Brand palette (GuideLearn)
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
PRIMARY = RGBColor(0x1E, 0x4F, 0x9E)
PRIMARY_SOFT = RGBColor(0x2F, 0x6F, 0xC3)
ACCENT = RGBColor(0xE8, 0xA3, 0x17)  # warm gold accent
TEAL = RGBColor(0x0D, 0x8A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF7, 0xFB)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
LIGHT_TEAL = RGBColor(0xE6, 0xF6, 0xF3)
LIGHT_GOLD = RGBColor(0xFF, 0xF6, 0xE0)
LIGHT_CORAL = RGBColor(0xFF, 0xEE, 0xEA)

OUT = Path(__file__).resolve().parents[1] / "GuideLearn-JUMP-THAILAND-2026.pptx"

W = Inches(13.333)
H = Inches(7.5)


def set_run(run, *, size=18, bold=False, color=SLATE, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    shape,
    text,
    *,
    size=18,
    bold=False,
    color=SLATE,
    align=PP_ALIGN.LEFT,
    font="Calibri",
    clear=True,
):
    tf = shape.text_frame
    if clear:
        tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(
    tf,
    text,
    *,
    size=16,
    bold=False,
    color=SLATE,
    align=PP_ALIGN.LEFT,
    space_before=6,
    space_after=0,
    font="Calibri",
):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return p


def rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    return s


def round_rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    # softer corners
    try:
        s.adjustments[0] = 0.12
    except Exception:
        pass
    return s


def footer(slide, page: int, total: int = 10):
    bar = rect(slide, 0, H - Inches(0.42), W, Inches(0.42), NAVY)
    box = slide.shapes.add_textbox(Inches(0.5), H - Inches(0.38), Inches(8), Inches(0.32))
    add_text(
        box,
        "GuideLearn  ·  JUMP THAILAND 2026  ·  AIS Ecosystem",
        size=11,
        color=RGBColor(0xB8, 0xC7, 0xDC),
    )
    num = slide.shapes.add_textbox(W - Inches(1.4), H - Inches(0.38), Inches(1.0), Inches(0.32))
    add_text(num, f"{page:02d} / {total:02d}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def section_label(slide, text, left=Inches(0.55), top=Inches(0.35)):
    pill = round_rect(slide, left, top, Inches(2.6), Inches(0.36), LIGHT_BLUE)
    add_text(pill, text, size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for p in pill.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER


def title_block(slide, title, subtitle=None):
    section_y = Inches(0.85)
    t = slide.shapes.add_textbox(Inches(0.55), section_y, Inches(12.2), Inches(0.7))
    add_text(t, title, size=32, bold=True, color=NAVY, font="Calibri")
    accent = rect(slide, Inches(0.55), Inches(1.5), Inches(1.15), Inches(0.08), ACCENT)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.4))
        add_text(s, subtitle, size=14, color=MUTED)


def card(slide, left, top, width, height, fill=CARD):
    # soft shadow plate
    shadow = round_rect(
        slide,
        left + Inches(0.04),
        top + Inches(0.05),
        width,
        height,
        RGBColor(0xD5, 0xDE, 0xEA),
    )
    body = round_rect(slide, left, top, width, height, fill)
    return body


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def slide_cover(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, NAVY)
    # decorative panels
    rect(slide, W - Inches(4.2), 0, Inches(4.2), H, PRIMARY)
    # diagonal-ish accent strip
    rect(slide, W - Inches(4.35), 0, Inches(0.18), H, ACCENT)

    brand = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(8), Inches(0.4))
    add_text(brand, "GUIDELEARN", size=16, bold=True, color=ACCENT)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(8.2), Inches(2.2))
    tf = add_text(
        title,
        "GuideLearn – AI Learning Coach",
        size=36,
        bold=True,
        color=WHITE,
        font="Calibri",
    )
    add_para(tf, "เพื่อความเท่าเทียมทางการศึกษา", size=30, bold=True, color=WHITE, space_before=8)

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(8.0), Inches(1.0))
    tf = add_text(
        sub,
        "ยกระดับการเรียนรู้เด็กไทยด้วย AI และ AIS Ecosystem",
        size=18,
        color=RGBColor(0xC9, 0xD8, 0xEC),
    )
    add_para(tf, "JUMP THAILAND 2026", size=16, bold=True, color=ACCENT, space_before=10)

    badge = round_rect(slide, W - Inches(3.55), Inches(2.7), Inches(2.9), Inches(2.0), WHITE)
    btf = add_text(badge, "Pitch Deck", size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_para(btf, "10 Slides", size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER, space_before=8)
    add_para(btf, "Prototype Ready", size=13, color=MUTED, align=PP_ALIGN.CENTER, space_before=6)

    foot = slide.shapes.add_textbox(Inches(0.7), H - Inches(0.85), Inches(8), Inches(0.35))
    add_text(foot, "Editable PowerPoint  ·  Competition Submission", size=12, color=RGBColor(0x8A, 0xA0, 0xBC))


def slide_problem(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "01  PROBLEM")
    title_block(slide, "ปัญหาและความจำเป็น", "ช่องว่างทางการศึกษาที่ต้องเร่งแก้")

    left = card(slide, Inches(0.55), Inches(2.35), Inches(5.9), Inches(3.9), LIGHT_CORAL)
    right = card(slide, Inches(6.85), Inches(2.35), Inches(5.9), Inches(3.9), LIGHT_GOLD)

    # left content
    icon = slide.shapes.add_textbox(Inches(0.85), Inches(2.6), Inches(5.3), Inches(0.4))
    add_text(icon, "ช่องว่างทางการศึกษา", size=20, bold=True, color=NAVY)
    body = slide.shapes.add_textbox(Inches(0.85), Inches(3.2), Inches(5.3), Inches(2.6))
    tf = add_text(
        body,
        "เด็กขาดแคลนเข้าไม่ถึงการติวเตอร์และสื่อการเรียนที่มีคุณภาพ",
        size=18,
        color=SLATE,
    )
    add_para(
        tf,
        "โอกาสเรียนรู้ยังกระจุกตัวในเมืองและครอบครัวที่มีกำลังซื้อ ทำให้เด็กในพื้นที่ห่างไกลเสียโอกาสแข่งขันตั้งแต่ต้นทาง",
        size=15,
        color=MUTED,
        space_before=14,
    )

    icon2 = slide.shapes.add_textbox(Inches(7.15), Inches(2.6), Inches(5.3), Inches(0.4))
    add_text(icon2, "ข้อจำกัดเชิงโครงสร้าง", size=20, bold=True, color=NAVY)
    body2 = slide.shapes.add_textbox(Inches(7.15), Inches(3.2), Inches(5.3), Inches(2.6))
    tf2 = add_text(body2, "งบประมาณน้อย", size=18, bold=True, color=SLATE)
    add_para(tf2, "ค่าติวและสื่อคุณภาพสูงยังเป็นภาระหนักของครัวเรือน", size=15, color=MUTED, space_before=6)
    add_para(tf2, "ปัญหาอินเทอร์เน็ตในพื้นที่ห่างไกล", size=18, bold=True, color=SLATE, space_before=16)
    add_para(
        tf2,
        "เน็ตช้า/ไม่เสถียรทำให้แพลตฟอร์มทั่วไปใช้งานไม่ได้จริงในบริบทโรงเรียนชายขอบ",
        size=15,
        color=MUTED,
        space_before=6,
    )
    footer(slide, 2)


def slide_solution(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "02  SOLUTION")
    title_block(slide, "ทางออก (Solution)", "AI Learning Coach ที่เข้าถึงได้จริง")

    items = [
        (
            "AI Personal Tutor",
            "ช่วยอธิบายโจทย์ทีละ Step ผ่านภาพถ่ายและเสียง เน้นแนะแนวคิดแบบโซเครติก ไม่เฉลยทันที",
            LIGHT_BLUE,
            PRIMARY,
        ),
        (
            "Accessible for All",
            "รองรับโหมดเน็ตช้า (Low-bandwidth) และออกแบบเพื่อความเท่าเทียม—ใช้ได้ทั้งมือถือและพื้นที่ห่างไกล",
            LIGHT_TEAL,
            TEAL,
        ),
    ]
    for i, (title, desc, bg, accent) in enumerate(items):
        top = Inches(2.35) + Inches(i * 2.15)
        c = card(slide, Inches(0.55), top, Inches(12.2), Inches(1.95), bg)
        bar = rect(slide, Inches(0.55), top, Inches(0.16), Inches(1.95), accent)
        t = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.35), Inches(11.3), Inches(0.45))
        add_text(t, title, size=22, bold=True, color=NAVY)
        d = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.9), Inches(11.3), Inches(0.8))
        add_text(d, desc, size=16, color=SLATE)
    footer(slide, 3)


def slide_features(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "03  KEY FEATURES")
    title_block(slide, "ฟีเจอร์หลัก", "เครื่องมือครบวงจรสำหรับนักเรียนและครู")

    features = [
        ("01", "Snap & Learn", "ถ่ายรูปโจทย์ แล้วให้ AI ช่วยแนะแนวคิดทีละขั้นตอน", LIGHT_BLUE),
        ("02", "Voice Interactive", "สั่งการและฟังคำอธิบายด้วยเสียง รองรับการเข้าถึง", LIGHT_TEAL),
        ("03", "Teacher Dashboard", "ช่วยครูติดตามเด็กกลุ่มเสี่ยงเรียนไม่ทันแบบเรียลไทม์", LIGHT_GOLD),
    ]
    width = Inches(3.9)
    gap = Inches(0.25)
    start = Inches(0.55)
    for i, (num, title, desc, bg) in enumerate(features):
        left = start + (width + gap) * i
        c = card(slide, left, Inches(2.4), width, Inches(3.85), bg)
        n = slide.shapes.add_textbox(left + Inches(0.35), Inches(2.7), width - Inches(0.7), Inches(0.5))
        add_text(n, num, size=28, bold=True, color=PRIMARY)
        t = slide.shapes.add_textbox(left + Inches(0.35), Inches(3.4), width - Inches(0.7), Inches(0.7))
        add_text(t, title, size=20, bold=True, color=NAVY)
        d = slide.shapes.add_textbox(left + Inches(0.35), Inches(4.25), width - Inches(0.7), Inches(1.6))
        add_text(d, desc, size=15, color=SLATE)
    footer(slide, 4)


def slide_tech(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "04  TECH ARCHITECTURE")
    title_block(slide, "สถาปัตยกรรมเทคโนโลยี", "Stack ที่พร้อมสเกลและดูแลรักษา")

    layers = [
        ("Frontend", "Next.js + TypeScript\nMobile-First UI ที่ใช้งานง่ายบนมือถือ", PRIMARY),
        ("AI Engine", "Vision Engine อ่านโจทย์จากภาพ\n+ Speech API สั่งงาน/ฟังคำอธิบาย", TEAL),
        ("Cloud Ready", "Deploy บน Vercel / พร้อมเชื่อม\nAIS Cloud ตามมาตรฐานความปลอดภัย", ACCENT),
    ]
    for i, (title, desc, color) in enumerate(layers):
        left = Inches(0.55) + Inches(i * 4.15)
        c = card(slide, left, Inches(2.45), Inches(3.95), Inches(3.7), WHITE)
        topbar = rect(slide, left, Inches(2.45), Inches(3.95), Inches(0.14), color)
        t = slide.shapes.add_textbox(left + Inches(0.35), Inches(2.85), Inches(3.25), Inches(0.5))
        add_text(t, title, size=20, bold=True, color=NAVY)
        d = slide.shapes.add_textbox(left + Inches(0.35), Inches(3.5), Inches(3.25), Inches(2.2))
        tf = d.text_frame
        tf.clear()
        for j, line in enumerate(desc.split("\n")):
            if j == 0:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = line
                set_run(run, size=15, color=SLATE)
            else:
                add_para(tf, line, size=15, color=SLATE, space_before=8)
    footer(slide, 5)


def slide_ais(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "05  AIS ECOSYSTEM")
    title_block(slide, "การเชื่อมต่อ AIS Ecosystem", "ต่อยอดโครงสร้างดิจิทัลของประเทศ")

    rows = [
        ("AIS Cloud", "THAI Hyperscale Cloud ปลอดภัยสูง จัดเก็บข้อมูลตามมาตรฐาน", LIGHT_BLUE),
        ("LearnDi", "แนะนำคอร์สเรียนออนไลน์ต่อยอดทักษะอนาคตให้ผู้เรียน", LIGHT_TEAL),
        ("Competency Assessment", "ระบบประเมินและค้นหาสมรรถนะผู้เรียน เพื่อวางแผนการเรียนรู้เฉพาะบุคคล", LIGHT_GOLD),
    ]
    for i, (title, desc, bg) in enumerate(rows):
        top = Inches(2.35) + Inches(i * 1.35)
        c = card(slide, Inches(0.55), top, Inches(12.2), Inches(1.2), bg)
        num = slide.shapes.add_textbox(Inches(0.85), top + Inches(0.35), Inches(0.6), Inches(0.5))
        add_text(num, f"{i+1:02d}", size=20, bold=True, color=PRIMARY)
        t = slide.shapes.add_textbox(Inches(1.6), top + Inches(0.22), Inches(10.6), Inches(0.4))
        add_text(t, title, size=18, bold=True, color=NAVY)
        d = slide.shapes.add_textbox(Inches(1.6), top + Inches(0.62), Inches(10.6), Inches(0.4))
        add_text(d, desc, size=14, color=SLATE)
    footer(slide, 6)


def slide_impact(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "06  IMPACT")
    title_block(slide, "ผลลัพธ์และผลกระทบ", "จากต้นแบบสู่การเปลี่ยนโอกาสเรียนรู้")

    left = card(slide, Inches(0.55), Inches(2.4), Inches(6.0), Inches(3.8), LIGHT_BLUE)
    right = card(slide, Inches(6.8), Inches(2.4), Inches(6.0), Inches(3.8), LIGHT_TEAL)

    t1 = slide.shapes.add_textbox(Inches(0.95), Inches(2.7), Inches(5.2), Inches(0.45))
    add_text(t1, "ระยะสั้น", size=14, bold=True, color=PRIMARY)
    h1 = slide.shapes.add_textbox(Inches(0.95), Inches(3.15), Inches(5.2), Inches(0.7))
    add_text(h1, "Web App Prototype", size=24, bold=True, color=NAVY)
    b1 = slide.shapes.add_textbox(Inches(0.95), Inches(4.0), Inches(5.2), Inches(1.8))
    add_text(
        b1,
        "ใช้งานได้จริงทันทีบนมือถือ รองรับติวการบ้านด้วยภาพ เสียง และโหมดเน็ตช้า",
        size=16,
        color=SLATE,
    )

    t2 = slide.shapes.add_textbox(Inches(7.2), Inches(2.7), Inches(5.2), Inches(0.45))
    add_text(t2, "ระยะยาว", size=14, bold=True, color=TEAL)
    h2 = slide.shapes.add_textbox(Inches(7.2), Inches(3.15), Inches(5.2), Inches(0.7))
    add_text(h2, "ลดความเหลื่อมล้ำ", size=24, bold=True, color=NAVY)
    b2 = slide.shapes.add_textbox(Inches(7.2), Inches(4.0), Inches(5.2), Inches(1.8))
    tf = add_text(
        b2,
        "ลดช่องว่างความเหลื่อมล้ำในการสอบเข้า",
        size=16,
        color=SLATE,
    )
    add_para(tf, "ลดภาระงานครู ด้วยแดชบอร์ดติดตามกลุ่มเสี่ยงและสรุปความก้าวหน้า", size=16, color=SLATE, space_before=12)
    footer(slide, 7)


def slide_business(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "07  BUSINESS MODEL")
    title_block(slide, "รูปแบบความยั่งยืน", "สร้างรายได้โดยไม่ทิ้งเด็กขาดแคลน")

    models = [
        ("Freemium", "เด็กขาดแคลนใช้งานฟีเจอร์พื้นฐานฟรี เพื่อเข้าถึงติวเตอร์ AI อย่างทั่วถึง", LIGHT_BLUE),
        ("B2G / School Subscription", "แพ็กเกจ Dashboard และสถิติลึกสำหรับโรงเรียน / สพฐ.", LIGHT_GOLD),
        ("Partnership", "เชื่อมต่อคอร์สเรียนเชิงลึกกับคอร์สพันธมิตรและ LearnDi", LIGHT_TEAL),
    ]
    for i, (title, desc, bg) in enumerate(models):
        left = Inches(0.55) + Inches(i * 4.15)
        c = card(slide, left, Inches(2.45), Inches(3.95), Inches(3.7), bg)
        idx = slide.shapes.add_textbox(left + Inches(0.3), Inches(2.75), Inches(3.3), Inches(0.4))
        add_text(idx, f"MODEL 0{i+1}", size=12, bold=True, color=PRIMARY)
        t = slide.shapes.add_textbox(left + Inches(0.3), Inches(3.25), Inches(3.3), Inches(1.0))
        add_text(t, title, size=18, bold=True, color=NAVY)
        d = slide.shapes.add_textbox(left + Inches(0.3), Inches(4.4), Inches(3.3), Inches(1.4))
        add_text(d, desc, size=14, color=SLATE)
    footer(slide, 8)


def slide_partners(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, OFF_WHITE)
    section_label(slide, "08  PARTNERSHIPS")
    title_block(slide, "พันธมิตรความร่วมมือ", "เครือข่ายที่ทำให้สเกลได้จริง")

    partners = [
        ("ศธ. / สพฐ.", "ทดลองใช้งานในโรงเรียนนำร่อง และเชื่อมระบบติดตามผลผู้เรียน"),
        ("กสศ.", "สนับสนุนเด็กทุนและโรงเรียนในพื้นที่ห่างไกล ให้เข้าถึงเครื่องมืออย่างทั่วถึง"),
        ("AIS Academy", "สนับสนุนโครงข่าย 5G และสื่อการเรียนรู้ดิจิทัลในระบบนิเวศ AIS"),
    ]
    for i, (title, desc) in enumerate(partners):
        top = Inches(2.35) + Inches(i * 1.35)
        c = card(slide, Inches(0.55), top, Inches(12.2), Inches(1.2), WHITE)
        accent = rect(slide, Inches(0.55), top, Inches(0.16), Inches(1.2), [PRIMARY, TEAL, ACCENT][i])
        t = slide.shapes.add_textbox(Inches(1.05), top + Inches(0.22), Inches(11.2), Inches(0.4))
        add_text(t, title, size=18, bold=True, color=NAVY)
        d = slide.shapes.add_textbox(Inches(1.05), top + Inches(0.62), Inches(11.2), Inches(0.4))
        add_text(d, desc, size=14, color=SLATE)
    footer(slide, 9)


def slide_prototype(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, 0, 0, Inches(0.18), H, ACCENT)

    section = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(10), Inches(0.35))
    add_text(section, "09  PROTOTYPE & REFERENCES", size=13, bold=True, color=ACCENT)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(12), Inches(0.7))
    add_text(title, "ต้นแบบและลิงก์อ้างอิง", size=32, bold=True, color=WHITE)

    # Live demo card
    c1 = round_rect(slide, Inches(0.7), Inches(2.3), Inches(5.9), Inches(2.6), WHITE)
    l1 = slide.shapes.add_textbox(Inches(1.05), Inches(2.55), Inches(5.2), Inches(0.4))
    add_text(l1, "LIVE DEMO", size=12, bold=True, color=PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(1.05), Inches(3.05), Inches(5.2), Inches(0.55))
    add_text(t1, "Web App Prototype", size=22, bold=True, color=NAVY)
    u1 = slide.shapes.add_textbox(Inches(1.05), Inches(3.7), Inches(5.2), Inches(0.8))
    add_text(u1, "guide-learn-zeta.vercel.app", size=16, color=PRIMARY)

    # GitHub card
    c2 = round_rect(slide, Inches(6.85), Inches(2.3), Inches(5.9), Inches(2.6), WHITE)
    l2 = slide.shapes.add_textbox(Inches(7.2), Inches(2.55), Inches(5.2), Inches(0.4))
    add_text(l2, "SOURCE CODE", size=12, bold=True, color=TEAL)
    t2 = slide.shapes.add_textbox(Inches(7.2), Inches(3.05), Inches(5.2), Inches(0.55))
    add_text(t2, "GitHub Repository", size=22, bold=True, color=NAVY)
    u2 = slide.shapes.add_textbox(Inches(7.2), Inches(3.7), Inches(5.2), Inches(0.8))
    add_text(u2, "github.com/cashchrys4762/GuideLearn", size=14, color=PRIMARY)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(12), Inches(1.0))
    tf = add_text(
        note,
        "สแกน/เปิดลิงก์เพื่อทดลอง Live Demo ได้ทันทีบนมือถือ",
        size=16,
        color=RGBColor(0xC9, 0xD8, 0xEC),
    )
    add_para(
        tf,
        "GuideLearn  ·  AI Learning Coach for Educational Equity  ·  JUMP THAILAND 2026",
        size=13,
        color=ACCENT,
        space_before=10,
    )


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_features(prs)
    slide_tech(prs)
    slide_ais(prs)
    slide_impact(prs)
    slide_business(prs)
    slide_partners(prs)
    slide_prototype(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
